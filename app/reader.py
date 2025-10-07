#PDF, DOCX, XLSX, CSV, PPTX Reader and Data Extractor
#Created By: Deepak Singh 
#Created on: 05-Oct-2025
#Modified On: 07-Oct-2025
#Version: 1.0.0.0
#---------------------------------------------

from pathlib import Path
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

def read_file(file_path):
    ext = Path(file_path).suffix.lower()
    text = ""
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
        elif ext == ".pdf":
            pdf = PdfReader(file_path)
            text = "\n".join((p.extract_text() or "") for p in pdf.pages)
        elif ext in [".xlsx", ".xls"]:
            xls = pd.ExcelFile(file_path)
            txts = []
            for s in xls.sheet_names:
                df = xls.parse(s, dtype=str)
                txts.append("\n".join(df.fillna("").astype(str).values.flatten()))
            text = "\n".join(txts)
        elif ext == ".pptx":
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                slide_txt = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_txt.append(shape.text)
                slides.append("\n".join(slide_txt))
            text = "\n".join(slides)
        elif ext == ".csv":
            df = pd.read_csv(file_path, dtype=str, encoding="utf-8", errors="ignore")
            text = "\n".join(df.fillna("").astype(str).values.flatten())
        else:
            print(f"[WARN] Unsupported file: {file_path}")
    except Exception as e:
        print(f"[WARN] Could not read {file_path}: {e}")
    return text
